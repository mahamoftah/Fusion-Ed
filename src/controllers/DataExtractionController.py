import os
import re
import sys
from urllib.parse import urlparse
from docx import Document
from fastapi import HTTPException, status
from pptx import Presentation
from langchain_community.document_loaders import PyMuPDFLoader
import requests
from src.controllers.BaseController import BaseController



class DataExtractionController(BaseController): 
    def __init__(self, file_urls):
        super().__init__()
        self.file_urls = file_urls
        self.file_extensions = [file_url.split('.')[-1] for file_url in file_urls]
        self.file_contents = []


    async def load(self):

        for index, (file_url, extension) in enumerate(zip(self.file_urls, self.file_extensions)):

            if not (await self.is_valid_url(file_url)):
                self.message = "has invalid URL"
                self.file_contents.append(
                    {
                        "success": False,
                        "message": self.message,
                        "content": None,
                        "index": index
                    }
                )
                self.logger.info(f"File {file_url} has invalid URL")
                continue

            content = await self.get_file_content(file_url, extension)

            if content and content.strip() == '':
                self.message = "is empty"
                self.file_contents.append(
                    {
                        "success": False,
                        "message": self.message,
                        "content": None,
                        "index": index
                    }
                )
                self.logger.info(f"File {file_url} is empty")
                continue

            elif content:
                processed_data = await self.process(content)
                self.message = "is successfully uploaded"
                self.file_contents.append(
                    {
                        "success": True,
                        "message": self.message,
                        "content": processed_data,
                        "index": index
                    }
                )
                continue

            else:
                self.message = "has unsupported extension"
                self.file_contents.append(
                    {
                        "success": False,
                        "message": self.message,
                        "content": None,
                        "index": index
                    }
                )

        return self.file_contents



    async def is_valid_url(self, file_url):
        if urlparse(file_url).scheme in ("http", "https"):
            try:
                response = requests.get(file_url)
                return response.status_code == 200
            except requests.exceptions.RequestException:
                return False
        else:
            return os.path.isfile(file_url)



    async def get_file_content(self, file_url, extension):
        self.logger.info(f"Extension: {extension}")
        if extension == "pdf":
            return await self.load_pdf(file_url)
        
        elif extension in {"docx", "doc"}: 
            return await self.load_docx(file_url)
        
        elif extension in {"pptx", "ppt"}:
            return await self.load_pptx(file_url)
        
        elif extension in {"md", "txt", "json"}:
            return await self.load_txt(file_url)
        
        else:
            return None


    async def load_docx(self, file_url):
        try:

            text = ""
            document = Document(file_url)
            text = "\n".join([p.text for p in document.paragraphs])
            
            return text
        except Exception:
            raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error in processing Word File")
    
    
    async def load_txt(self, file_url):
        try:
            with open(file_url, "r", encoding="utf-8") as f:
                return f.read()
        except Exception:
            raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error in processing txt File")
    


    async def load_pptx(self, file_url):
        try:
            text = ""

            presentation = Presentation(file_url)
                
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text += paragraph.text + "\n"

            return text
        except Exception:
            raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error in processing powerpoint File")
        
    

    async def load_pdf(self, file_url):
        try:
            loader = PyMuPDFLoader(file_url)
            row_data = ""
            for doc in loader.load():
                row_data += doc.page_content + "\n"
            return row_data
        
        except Exception as e:
            self.logger.error(f"Error processing PDF file: {str(e)}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Error in processing PDF file: {str(e)}"
            )
    

    async def process(self, row_data):

        processed_data = re.sub(r' +', ' ', row_data)
        processed_data = re.sub(r'\n\s*\n', '\n', processed_data).strip()

        return processed_data